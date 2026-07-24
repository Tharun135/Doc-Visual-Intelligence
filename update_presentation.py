from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 45, 85)  # Dark blue
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 220, 255)

def add_content_slide(prs, title, content_text):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 245, 250)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 45, 85)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for line in content_text.strip().split('\n'):
        line = line.strip()
        if not line:
            p = text_frame.add_paragraph()
            p.text = ""
        elif line.startswith('- '):
            p = text_frame.add_paragraph()
            p.text = line[2:]
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(6)
        elif line.startswith('  - '):
            p = text_frame.add_paragraph()
            p.text = line[4:]
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(80, 80, 80)
        else:
            p = text_frame.add_paragraph() if text_frame.paragraphs[-1].text else text_frame.paragraphs[0]
            if p.text:
                p = text_frame.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(50, 50, 50)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1
add_title_slide(prs, "Doc Visual Intelligence", 
                "Intelligent Visual Recommendation Engine for Technical Documentation\n\nAnalyze. Recommend. Generate.")

# Slide 2
add_content_slide(prs, "The Problem", """
- Writers manually decide where visuals are needed
- Visual decisions vary between authors
- Important sections remain text-only
- Reviews focus on missing diagrams rather than content quality
- Adding visuals late increases documentation effort

Key Question: Which sections actually deserve a visual?
""")

# Slide 3
add_content_slide(prs, "Our Solution", """
Doc Visual Intelligence automatically analyzes every document section and determines:

- Does this section need a visual?
- Which visual communicates it best?
- Why is that visual recommended?
- How confident is the recommendation?
- Can the visual be generated automatically?
""")

# Slide 4
add_content_slide(prs, "How It Works", """
Supported Input:
- Markdown
- DOCX
- PDF
- AsciiDoc
- Plain Text

Output:
- Visual recommendation
- Confidence level
- Explanation
- SVG diagram
""")

# Slide 5
add_content_slide(prs, "Intelligence Pipeline", """
- Signal Extraction
- Content Classification
- Visual Worthiness Analysis
- Rule-Based Recommendation
- Confidence Scoring
- Visual Generation

Each section passes through every stage before a recommendation is produced.
""")

# Slide 6
add_content_slide(prs, "Signal Extraction", """
The engine converts document sections into measurable characteristics.

Signals include:
- Procedural steps
- UI interactions
- Component relationships
- Data flow
- Conditional branches
- Complexity metrics
- Word count
- Network entities
- Comparison patterns

These signals provide objective evidence for visual recommendations.
""")

# Slide 7
add_content_slide(prs, "Recommendation Engine", """
Content Classification:
- Procedure
- Architecture
- Troubleshooting
- Reference
- Concept

Recommended Visual Types:
- Workflow Diagram
- Screenshot
- Architecture Diagram
- Data Flow Diagram
- Decision Tree
- Sequence Diagram
- Comparison Table

Every recommendation includes an explainable confidence score.
""")

# Slide 8
add_content_slide(prs, "Recommendation Example", """
Input: "Open Industrial Edge Management. Select the device. Install the connector. Configure endpoint. Deploy."

Classification: Procedure
Signals: Steps = 5, UI Interactions = 4, Worthiness = 8/10

Recommendations:
- Workflow Diagram (Confidence 94%)
- Screenshot (Confidence 88%)

Reason: High procedural density, Multiple UI actions, Sequential execution
""")

# Slide 9
add_content_slide(prs, "Supported Visual Types", """
Procedural:
- Workflow Diagram, Flowchart, GIF Tutorial

User Interface:
- Screenshot, Configuration Screenshot

Architecture:
- Architecture Diagram, Topology Diagram, Data Flow Diagram, Sequence Diagram

Logic:
- Decision Tree, Comparison Table

Reference:
- Mapping Table, Code Example, Visual Summary
""")

# Slide 10
add_content_slide(prs, "Business Value", """
Benefits for Technical Writers:
- Consistent visual recommendations
- Faster documentation reviews
- Reduced manual planning
- Better reader comprehension
- Improved documentation quality

Benefits for Organizations:
- Standardized documentation
- Reduced support effort
- Faster onboarding
- Improved documentation consistency
""")

# Slide 11
add_content_slide(prs, "Future Roadmap", """
Near Term:
- Improve recommendation rules
- Domain-specific terminology
- Better confidence calibration

Mid Term:
- Learning from user feedback
- Smarter prioritization

Long Term:
- VS Code integration
- Markdown extension
- CI/CD documentation pipeline
- Enterprise documentation platforms
""")

# Slide 12
add_content_slide(prs, "Thank You", """
Helping Technical Writers Build Visual-First Documentation

Questions?
""")

# Save
prs.save('doc-agent-orchestrator.pptx')
print("✓ Presentation updated successfully with 12 slides!")
