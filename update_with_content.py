from pptx import Presentation
from pptx.util import Inches, Pt

# Load the original presentation to preserve styling
prs = Presentation('Doc-Visual-Intelligence.pptx')

# Clear existing slides but keep the master/theme
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# Define slide content
slide_contents = [
    {
        "title": "Doc Visual Intelligence",
        "content": "Intelligent Visual Recommendation Engine for Technical Documentation\n\nAnalyze. Recommend. Generate.",
    },
    {
        "title": "The Problem",
        "content": "Technical documentation is often text-heavy, making complex procedures and architectures difficult to understand.\n\nCurrent challenges:\n• Writers manually decide where visuals are needed\n• Visual decisions vary between authors\n• Important sections remain text-only\n• Reviews focus on missing diagrams rather than content quality\n• Adding visuals late increases documentation effort\n\nKey Question: Which sections actually deserve a visual?",
    },
    {
        "title": "Our Solution",
        "content": "Doc Visual Intelligence automatically analyzes every document section and determines:\n• Does this section need a visual?\n• Which visual communicates it best?\n• Why is that visual recommended?\n• How confident is the recommendation?\n• Can the visual be generated automatically?",
    },
    {
        "title": "How It Works",
        "content": "Supported input:\n• Markdown\n• DOCX\n• PDF\n• AsciiDoc\n• Plain Text\n\nOutput:\n• Visual recommendation\n• Confidence\n• Explanation\n• SVG diagram",
    },
    {
        "title": "Intelligence Pipeline",
        "content": "• Signal Extraction\n• Content Classification\n• Visual Worthiness\n• Rule-Based Recommendation\n• Confidence Scoring\n• Visual Generation\n\nEach section passes through every stage before a recommendation is produced.",
    },
    {
        "title": "Signal Extraction",
        "content": "The engine converts document sections into measurable characteristics.\n\nSignals include:\n• Procedural steps\n• UI interactions\n• Component relationships\n• Data flow\n• Conditional branches\n• Complexity\n• Word count\n• Network entities\n• Comparison patterns\n\nThese signals provide objective evidence for visual recommendations.",
    },
    {
        "title": "Recommendation Engine",
        "content": "Each section is classified as:\n• Procedure\n• Architecture\n• Troubleshooting\n• Reference\n• Concept\n\nThe engine then evaluates specialized rules to recommend visuals such as:\n• Workflow Diagram\n• Screenshot\n• Architecture Diagram\n• Data Flow Diagram\n• Decision Tree\n• Sequence Diagram\n• Comparison Table\n\nEvery recommendation includes an explainable confidence score.",
    },
    {
        "title": "Recommendation Example",
        "content": "Input: Open Industrial Edge Management. Select the device. Install the connector. Configure endpoint. Deploy.\n\nClassification: Procedure\nSignals: Steps = 5, UI Interactions = 4, Worthiness = 8/10\n\nRecommendation:\n• Workflow Diagram (Confidence 94%)\n• Screenshot (Confidence 88%)\n\nReason: High procedural density, Multiple UI actions, Sequential execution",
    },
    {
        "title": "Supported Visual Types",
        "content": "Procedural: Workflow Diagram, Flowchart, GIF Tutorial\n\nUser Interface: Screenshot, Configuration Screenshot\n\nArchitecture: Architecture Diagram, Topology Diagram, Data Flow Diagram, Sequence Diagram\n\nLogic: Decision Tree, Comparison Table\n\nReference: Mapping Table, Code Example, Visual Summary",
    },
    {
        "title": "Business Value",
        "content": "Benefits for Technical Writers:\n• Consistent visual recommendations\n• Faster documentation reviews\n• Reduced manual planning\n• Better reader comprehension\n• Improved documentation quality\n\nBenefits for Organizations:\n• Standardized documentation\n• Reduced support effort\n• Faster onboarding\n• Improved documentation consistency",
    },
    {
        "title": "Future Roadmap",
        "content": "Near Term:\n• Improve recommendation rules\n• Domain-specific terminology\n• Better confidence calibration\n\nMid Term:\n• Learning from user feedback\n• Smarter prioritization\n\nLong Term:\n• VS Code integration\n• Markdown extension\n• CI/CD documentation pipeline\n• Enterprise documentation platforms",
    },
    {
        "title": "Thank You",
        "content": "Helping Technical Writers Build Visual-First Documentation\n\nQuestions?",
    },
]

# Add slides with content
for slide_data in slide_contents:
    # Use blank layout to start fresh
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    content_text = slide_data["content"]
    lines = content_text.split('\n')
    
    for idx, line in enumerate(lines):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.font.size = Pt(18)
        p.level = 0

prs.save('Doc-Visual-Intelligence.pptx')
print("✓ Presentation updated successfully!")
print(f"✓ Total slides: {len(prs.slides)}")
