from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation('Doc-Visual-Intelligence.pptx')
print(f'Total slides: {len(prs.slides)}')
print(f'Slide dimensions: {prs.slide_width.inches}" x {prs.slide_height.inches}"')
print()

for i, slide in enumerate(prs.slides, 1):
    print(f'Slide {i}:')
    print(f'  Layout: {slide.slide_layout.name}')
    print(f'  Shapes: {len(slide.shapes)}')
    for j, shape in enumerate(slide.shapes):
        print(f'    Shape {j}:')
        print(f'      Type: {shape.shape_type}')
        if hasattr(shape, 'text'):
            text_preview = shape.text[:100].replace('\n', ' ')
            print(f'      Text: {text_preview}')
        if hasattr(shape, 'fill'):
            print(f'      Fill type: {shape.fill.type}')
        if hasattr(shape, 'left'):
            print(f'      Position: ({shape.left.inches:.2f}", {shape.top.inches:.2f}")')
            print(f'      Size: {shape.width.inches:.2f}" x {shape.height.inches:.2f}"')
    print()
