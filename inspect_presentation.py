from pptx import Presentation

prs = Presentation('doc-agent-orchestrator.pptx')
print(f'Total slides: {len(prs.slides)}')
print(f'Slide dimensions: {prs.slide_width.inches}" x {prs.slide_height.inches}"')
print()

for i, slide in enumerate(prs.slides, 1):
    print(f'Slide {i}:')
    print(f'  Layout: {slide.slide_layout.name}')
    print(f'  Shapes: {len(slide.shapes)}')
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text_preview = shape.text[:80].replace('\n', ' ')
            print(f'    - {shape.shape_type}: {text_preview}...' if len(shape.text) > 80 else f'    - {shape.shape_type}: {text_preview}')
    print()
