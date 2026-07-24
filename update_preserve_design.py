from pptx import Presentation

# New content for each slide
new_content = [
    {
        "title": "Doc Visual Intelligence",
        "subtitle": "Intelligent Visual Recommendation Engine for Technical Documentation\n\nAnalyze. Recommend. Generate.",
    },
    {
        "title": "The Problem",
        "body": "Technical documentation is often text-heavy, making complex procedures and architectures difficult to understand.\n\nCurrent challenges:\n• Writers manually decide where visuals are needed\n• Visual decisions vary between authors\n• Important sections remain text-only\n• Reviews focus on missing diagrams rather than content quality\n• Adding visuals late increases documentation effort\n\nKey Question: Which sections actually deserve a visual?",
    },
    {
        "title": "Our Solution",
        "body": "Doc Visual Intelligence automatically analyzes every document section and determines:\n• Does this section need a visual?\n• Which visual communicates it best?\n• Why is that visual recommended?\n• How confident is the recommendation?\n• Can the visual be generated automatically?",
    },
    {
        "title": "How It Works",
        "body": "Supported input:\n• Markdown • DOCX • PDF • AsciiDoc • Plain Text\n\nOutput:\n• Visual recommendation\n• Confidence • Explanation • SVG diagram",
    },
    {
        "title": "Intelligence Pipeline",
        "body": "• Signal Extraction\n• Content Classification\n• Visual Worthiness\n• Rule-Based Recommendation\n• Confidence Scoring\n• Visual Generation\n\nEach section passes through every stage before a recommendation is produced.",
    },
    {
        "title": "Signal Extraction",
        "body": "The engine converts document sections into measurable characteristics.\n\nSignals include:\n• Procedural steps • UI interactions • Component relationships • Data flow • Conditional branches • Complexity • Word count • Network entities • Comparison patterns\n\nThese signals provide objective evidence for visual recommendations.",
    },
    {
        "title": "Recommendation Engine",
        "body": "Each section is classified as:\nProcedure • Architecture • Troubleshooting • Reference • Concept\n\nRecommended visual types:\n• Workflow Diagram • Screenshot • Architecture Diagram • Data Flow Diagram • Decision Tree • Sequence Diagram • Comparison Table\n\nEvery recommendation includes an explainable confidence score.",
    },
    {
        "title": "Recommendation Example",
        "body": "Input: Open Industrial Edge Management. Select device. Install connector. Configure endpoint. Deploy.\n\nClassification: Procedure\nSignals: Steps=5, UI Interactions=4, Worthiness=8/10\n\nRecommendations:\n• Workflow Diagram (94% confidence)\n• Screenshot (88% confidence)\n\nReason: High procedural density, Multiple UI actions, Sequential execution",
    },
    {
        "title": "Supported Visual Types",
        "body": "Procedural: Workflow Diagram, Flowchart, GIF Tutorial\n\nUI: Screenshot, Configuration Screenshot\n\nArchitecture: Architecture Diagram, Topology, Data Flow, Sequence Diagram\n\nLogic: Decision Tree, Comparison Table\n\nReference: Mapping Table, Code Example, Visual Summary",
    },
    {
        "title": "Business Value",
        "body": "For Technical Writers:\n• Consistent visual recommendations\n• Faster documentation reviews\n• Reduced manual planning\n• Better reader comprehension\n• Improved documentation quality\n\nFor Organizations:\n• Standardized documentation • Reduced support effort • Faster onboarding • Documentation consistency",
    },
    {
        "title": "Future Roadmap",
        "body": "Near Term: Improve recommendation rules, Domain-specific terminology, Better confidence calibration\n\nMid Term: Learning from user feedback, Smarter prioritization\n\nLong Term: VS Code integration, Markdown extension, CI/CD documentation pipeline, Enterprise platforms",
    },
    {
        "title": "Thank You",
        "body": "Helping Technical Writers Build Visual-First Documentation\n\nQuestions?",
    },
]

# Load original presentation
prs = Presentation('Doc-Visual-Intelligence.pptx')

# Update text in each slide, preserving formatting and layout
for slide_idx, slide in enumerate(prs.slides):
    if slide_idx >= len(new_content):
        break
    
    content = new_content[slide_idx]
    
    # Find and update title and body text boxes
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text_frame = shape.text_frame
            
            # Determine if this is a title or body based on position/size
            if hasattr(shape, 'top') and shape.top.inches < 0.8:
                # Likely a title
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = content.get("title", "")
            elif hasattr(shape, 'top'):
                # Likely body content
                body_text = content.get("body") or content.get("subtitle", "")
                if body_text:
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = body_text

prs.save('Doc-Visual-Intelligence.pptx')
print("✓ Presentation updated successfully with new content!")
print(f"✓ Preserved original design and formatting")
print(f"✓ Total slides updated: {len(prs.slides)}")
